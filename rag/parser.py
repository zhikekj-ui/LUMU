"""Document parser — supports PDF, Word, Excel, PPT, TXT, Markdown, CSV, HTML."""
import os
from pathlib import Path
from typing import Optional


class DocumentChunk:
    """A chunk of text extracted from a document."""

    def __init__(self, text: str, metadata: dict):
        self.text = text
        self.metadata = metadata  # source, page, section, etc.

    def __repr__(self):
        src = self.metadata.get("source", "?")
        return f"Chunk({src}, {len(self.text)} chars)"


class DocumentParser:
    """Parse documents into text chunks for RAG indexing."""

    SUPPORTED_EXTENSIONS = {
        ".pdf", ".docx", ".doc", ".xlsx", ".xls", ".pptx",
        ".txt", ".md", ".csv", ".html", ".htm", ".json",
    }

    def __init__(self, chunk_size: int = 500, chunk_overlap: int = 50):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap

    def parse_file(self, file_path: str) -> list[DocumentChunk]:
        """Parse a single file into chunks."""
        path = Path(file_path)
        ext = path.suffix.lower()
        metadata_base = {"source": str(path), "filename": path.name}

        if ext == ".pdf":
            text = self._parse_pdf(path)
        elif ext in (".docx", ".doc"):
            text = self._parse_docx(path)
        elif ext in (".xlsx", ".xls"):
            text = self._parse_excel(path)
        elif ext == ".pptx":
            text = self._parse_pptx(path)
        elif ext in (".txt", ".md"):
            text = path.read_text(encoding="utf-8", errors="replace")
        elif ext == ".csv":
            text = self._parse_csv(path)
        elif ext in (".html", ".htm"):
            text = self._parse_html(path)
        elif ext == ".json":
            text = path.read_text(encoding="utf-8", errors="replace")
        else:
            raise ValueError(f"Unsupported file type: {ext}")

        return self._chunk_text(text, metadata_base)

    def parse_text(self, text: str, source: str = "direct_input") -> list[DocumentChunk]:
        """Parse raw text into chunks."""
        return self._chunk_text(text, {"source": source, "filename": source})

    def _parse_pdf(self, path: Path) -> str:
        """Parse PDF using PyPDF2."""
        try:
            from PyPDF2 import PdfReader
        except ImportError:
            raise ImportError("PyPDF2 not installed. Run: pip install PyPDF2")
        reader = PdfReader(str(path))
        pages = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text:
                pages.append(text)
        return "\n\n".join(pages)

    def _parse_docx(self, path: Path) -> str:
        """Parse Word document using python-docx."""
        try:
            from docx import Document
        except ImportError:
            raise ImportError("python-docx not installed. Run: pip install python-docx")
        doc = Document(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paragraphs)

    def _parse_excel(self, path: Path) -> str:
        """Parse Excel using openpyxl."""
        try:
            from openpyxl import load_workbook
        except ImportError:
            raise ImportError("openpyxl not installed. Run: pip install openpyxl")
        wb = load_workbook(str(path), read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    rows.append(" | ".join(cells))
            if rows:
                parts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(parts)

    def _parse_pptx(self, path: Path) -> str:
        """Parse PowerPoint using python-pptx."""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx not installed. Run: pip install python-pptx")
        prs = Presentation(str(path))
        parts = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                parts.append(f"[Slide {i + 1}]\n" + "\n".join(texts))
        return "\n\n".join(parts)

    def _parse_csv(self, path: Path) -> str:
        """Parse CSV file."""
        return path.read_text(encoding="utf-8", errors="replace")

    def _parse_html(self, path: Path) -> str:
        """Parse HTML, stripping tags."""
        try:
            from html.parser import HTMLParser

            class TextExtractor(HTMLParser):
                def __init__(self):
                    super().__init__()
                    self.text_parts = []

                def handle_data(self, data):
                    self.text_parts.append(data)

            parser = TextExtractor()
            html = path.read_text(encoding="utf-8", errors="replace")
            parser.feed(html)
            return "\n".join(parser.text_parts)
        except Exception:
            return path.read_text(encoding="utf-8", errors="replace")

    def _chunk_text(self, text: str, metadata: dict) -> list[DocumentChunk]:
        """Split text into overlapping chunks."""
        if not text.strip():
            return []

        chunks = []
        start = 0
        text_len = len(text)
        chunk_idx = 0

        while start < text_len:
            end = start + self.chunk_size
            chunk_text = text[start:end].strip()

            if chunk_text:
                chunk_meta = {**metadata, "chunk_index": chunk_idx, "char_offset": start}
                chunks.append(DocumentChunk(chunk_text, chunk_meta))
                chunk_idx += 1

            start = end - self.chunk_overlap
            if start >= text_len:
                break

        return chunks

    def supported(self, file_path: str) -> bool:
        """Check if a file type is supported."""
        ext = Path(file_path).suffix.lower()
        return ext in self.SUPPORTED_EXTENSIONS
